"""Office document parser — DOCX, XLSX, PPTX.

Uses ``python-docx``, ``openpyxl``, and ``python-pptx`` when available.
Each is an optional dependency; the parser gracefully skips formats whose
library is missing.
"""

from __future__ import annotations

import io
import logging

from src.core.models import FileRecord, ParseResult
from src.parsers.base import ParserBase

logger = logging.getLogger(__name__)

_OFFICE_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


class OfficeParser(ParserBase):
    name = "office"
    supported_mimes = list(_OFFICE_MIMES)
    lane = "heavy"

    async def can_parse(self, file_record: FileRecord) -> bool:
        if file_record.mime_type in _OFFICE_MIMES:
            return True
        from pathlib import Path

        return Path(file_record.identity.path).suffix.lower() in _OFFICE_EXTENSIONS

    async def parse(self, file_record: FileRecord, raw_bytes: bytes) -> ParseResult:
        from pathlib import Path

        ext = Path(file_record.identity.path).suffix.lower()

        if ext == ".docx":
            return self._parse_docx(raw_bytes)
        elif ext == ".xlsx":
            return self._parse_xlsx(raw_bytes)
        elif ext == ".pptx":
            return self._parse_pptx(raw_bytes)

        return ParseResult(content="", metadata={"error": "unknown_office_format"})

    @staticmethod
    def _parse_docx(data: bytes) -> ParseResult:
        try:
            from docx import Document  # type: ignore[import-untyped]
        except ImportError:
            return ParseResult(content="", metadata={"error": "python-docx not installed"})

        try:
            doc = Document(io.BytesIO(data))
        except Exception as exc:
            return ParseResult(content="", metadata={"error": f"corrupt_docx: {exc}"})

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n\n".join(paragraphs)
        return ParseResult(
            content=content,
            metadata={"paragraph_count": len(paragraphs), "char_count": len(content)},
        )

    @staticmethod
    def _parse_xlsx(data: bytes) -> ParseResult:
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
        except ImportError:
            return ParseResult(content="", metadata={"error": "openpyxl not installed"})

        try:
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        except Exception as exc:
            return ParseResult(content="", metadata={"error": f"corrupt_xlsx: {exc}"})

        try:
            rows: list[str] = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows.append(f"## Sheet: {sheet}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows.append("\t".join(cells))
        finally:
            wb.close()

        content = "\n".join(rows)
        return ParseResult(
            content=content,
            metadata={"sheet_count": len(wb.sheetnames), "char_count": len(content)},
        )

    @staticmethod
    def _parse_pptx(data: bytes) -> ParseResult:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError:
            return ParseResult(content="", metadata={"error": "python-pptx not installed"})

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            return ParseResult(content="", metadata={"error": f"corrupt_pptx: {exc}"})
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
            slides_text.append("\n".join(parts))

        content = "\n\n".join(slides_text)
        return ParseResult(
            content=content,
            metadata={"slide_count": len(prs.slides), "char_count": len(content)},
        )
